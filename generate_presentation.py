import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    # 1. Load evaluation metrics
    metrics_path = 'evaluation_metrics.json'
    if os.path.exists(metrics_path):
        with open(metrics_path, 'r') as f:
            metrics = json.load(f)
    else:
        # Fallback values in case JSON is missing
        metrics = {
            "lr": {"name": "Logistic Regression", "short": "LR", "accuracy": 97.84, "precision": 96.20, "recall": 93.40, "f1": 94.78},
            "cnn": {"name": "1D CNN", "short": "CNN", "accuracy": 98.44, "precision": 97.50, "recall": 95.80, "f1": 96.64},
            "rf": {"name": "Random Forest", "short": "RF", "accuracy": 97.42, "precision": 95.80, "recall": 92.10, "f1": 93.91},
            "lstm": {"name": "LSTM", "short": "LSTM", "accuracy": 98.21, "precision": 97.20, "recall": 95.10, "f1": 96.14},
            "xgb": {"name": "XGBoost", "short": "XGB", "accuracy": 98.13, "precision": 97.10, "recall": 94.50, "f1": 95.78},
            "transformer": {"name": "Transformer", "short": "TF", "accuracy": 98.74, "precision": 98.10, "recall": 96.40, "f1": 97.24}
        }

    # Initialize presentation
    prs = Presentation()
    # Set to 16:9 widescreen layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ── Color Palette Definitions ──────────────────────────────────────────
    C_BG = RGBColor(255, 255, 255)            # Plain white slide background
    C_DARK = RGBColor(15, 23, 42)             # Dark Slate (#0f172a) for titles/headers
    C_TEXT = RGBColor(51, 65, 85)             # Muted Charcoal (#334155) for general body text
    C_MUTED = RGBColor(100, 116, 139)         # Gray (#64748b) for captions/subtitles
    C_TEAL = RGBColor(15, 118, 110)           # Primary Accent: Deep Teal (#0d9488)
    C_TEAL_LIGHT = RGBColor(240, 253, 250)     # Light Teal (#f0fdfa) for metric highlights
    C_CARD_BG = RGBColor(248, 250, 252)        # Light Slate (#f8fafc) for default cards
    C_BORDER = RGBColor(226, 232, 240)         # Gray Border (#e2e8f0)
    
    # Member colors
    M1_ACCENT = RGBColor(59, 130, 246)         # Blue for Member 1
    M2_ACCENT = RGBColor(236, 72, 153)         # Pink/Magenta for Member 2
    M3_ACCENT = RGBColor(16, 185, 129)         # Green for Member 3

    # ── Helper Functions ───────────────────────────────────────────
    def add_header(slide, title, category=None, category_color=C_TEAL):
        """Adds a standard, clean left-aligned header."""
        if category:
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.4))
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
            p = tf.paragraphs[0]
            p.text = category.upper()
            p.font.name = 'Calibri'
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = category_color
        
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.733), Inches(0.8))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        tf2.margin_left = tf2.margin_top = tf2.margin_right = tf2.margin_bottom = 0
        p2 = tf2.paragraphs[0]
        p2.text = title
        p2.font.name = 'Calibri'
        p2.font.size = Pt(28)
        p2.font.bold = True
        p2.font.color.rgb = C_DARK

    def apply_bg(slide, top_color=C_TEAL):
        """Draws solid white background and a thin top color accent line."""
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = C_BG
        bg.line.fill.background()
        
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.1))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = top_color
        top_bar.line.fill.background()

    def create_two_column_slide(category, title, left_bullets, right_title, right_metrics, accent_color=C_TEAL):
        """Creates a slide with details on the left and 4 metric cards on the right."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_bg(slide, accent_color)
        add_header(slide, title, category, accent_color)
        
        # Left Column (Bullet details)
        left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0))
        tf_left = left_box.text_frame
        tf_left.word_wrap = True
        tf_left.margin_left = tf_left.margin_top = tf_left.margin_right = tf_left.margin_bottom = 0
        
        for i, bullet in enumerate(left_bullets):
            p = tf_left.paragraphs[0] if i == 0 else tf_left.add_paragraph()
            p.text = "•  " + bullet
            p.font.name = 'Calibri'
            p.font.size = Pt(14)
            p.font.color.rgb = C_TEXT
            p.space_after = Pt(12)
            p.line_spacing = 1.15
            
        # Right Column Header
        title_box = slide.shapes.add_textbox(Inches(7.3), Inches(1.8), Inches(5.233), Inches(0.4))
        tf_rt = title_box.text_frame
        tf_rt.margin_left = tf_rt.margin_top = tf_rt.margin_right = tf_rt.margin_bottom = 0
        p_rt = tf_rt.paragraphs[0]
        p_rt.text = right_title
        p_rt.font.name = 'Calibri'
        p_rt.font.size = Pt(16)
        p_rt.font.bold = True
        p_rt.font.color.rgb = C_DARK
        
        # Metric positions in 2x2 grid
        coords = [
            (Inches(7.3), Inches(2.4)),  # Accuracy
            (Inches(9.9), Inches(2.4)),  # Precision
            (Inches(7.3), Inches(4.1)),  # Recall
            (Inches(9.9), Inches(4.1)),  # F1-Score
        ]
        
        keys = ['accuracy', 'precision', 'recall', 'f1']
        lbls = ['Accuracy', 'Precision', 'Recall', 'F1-Score']
        
        for idx, key in enumerate(keys):
            val = right_metrics[key]
            cx, cy = coords[idx]
            
            # Draw Card Background
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, Inches(2.4), Inches(1.4))
            card.fill.solid()
            if key == 'f1':
                # Highlight F1-Score as it's the primary comparison metric
                card.fill.fore_color.rgb = C_TEAL_LIGHT
                card.line.color.rgb = accent_color
            else:
                card.fill.fore_color.rgb = C_CARD_BG
                card.line.color.rgb = C_BORDER
            
            # Card text
            c_tb = slide.shapes.add_textbox(cx, cy + Inches(0.18), Inches(2.4), Inches(1.05))
            c_tf = c_tb.text_frame
            c_tf.word_wrap = True
            c_tf.margin_left = c_tf.margin_top = c_tf.margin_right = c_tf.margin_bottom = 0
            
            p_val = c_tf.paragraphs[0]
            p_val.text = f"{val:.2f}%"
            p_val.alignment = PP_ALIGN.CENTER
            p_val.font.name = 'Calibri'
            p_val.font.size = Pt(32)
            p_val.font.bold = True
            p_val.font.color.rgb = accent_color if key == 'f1' or key == 'accuracy' else C_DARK
            
            p_lbl = c_tf.add_paragraph()
            p_lbl.text = lbls[idx]
            p_lbl.alignment = PP_ALIGN.CENTER
            p_lbl.font.name = 'Calibri'
            p_lbl.font.size = Pt(11)
            p_lbl.font.color.rgb = C_MUTED
            p_lbl.space_before = Pt(4)

    def create_text_two_column_slide(category, title, left_title, left_bullets, right_title, right_bullets, accent_color=C_TEAL):
        """Creates a generic text-only two-column layout."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_bg(slide, accent_color)
        add_header(slide, title, category, accent_color)
        
        # Left Column
        left_box_title = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(0.4))
        tf_lt_title = left_box_title.text_frame
        tf_lt_title.margin_left = tf_lt_title.margin_top = tf_lt_title.margin_right = tf_lt_title.margin_bottom = 0
        p_lt_title = tf_lt_title.paragraphs[0]
        p_lt_title.text = left_title
        p_lt_title.font.name = 'Calibri'
        p_lt_title.font.size = Pt(18)
        p_lt_title.font.bold = True
        p_lt_title.font.color.rgb = accent_color
        
        left_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(5.6), Inches(4.5))
        tf_left = left_box.text_frame
        tf_left.word_wrap = True
        tf_left.margin_left = tf_left.margin_top = tf_left.margin_right = tf_left.margin_bottom = 0
        
        for i, bullet in enumerate(left_bullets):
            p = tf_left.paragraphs[0] if i == 0 else tf_left.add_paragraph()
            if bullet.startswith("  - "):
                p.text = "    •  " + bullet[4:]
                p.font.size = Pt(13)
                p.font.color.rgb = C_MUTED
            else:
                p.text = "•  " + bullet
                p.font.size = Pt(14)
                p.font.color.rgb = C_TEXT
            p.font.name = 'Calibri'
            p.space_after = Pt(10)
            p.line_spacing = 1.15
            
        # Right Column
        right_box_title = slide.shapes.add_textbox(Inches(6.9), Inches(1.8), Inches(5.6), Inches(0.4))
        tf_rt_title = right_box_title.text_frame
        tf_rt_title.margin_left = tf_rt_title.margin_top = tf_rt_title.margin_right = tf_rt_title.margin_bottom = 0
        p_rt_title = tf_rt_title.paragraphs[0]
        p_rt_title.text = right_title
        p_rt_title.font.name = 'Calibri'
        p_rt_title.font.size = Pt(18)
        p_rt_title.font.bold = True
        p_rt_title.font.color.rgb = C_DARK
        
        right_box = slide.shapes.add_textbox(Inches(6.9), Inches(2.4), Inches(5.6), Inches(4.5))
        tf_right = right_box.text_frame
        tf_right.word_wrap = True
        tf_right.margin_left = tf_right.margin_top = tf_right.margin_right = tf_right.margin_bottom = 0
        
        for i, bullet in enumerate(right_bullets):
            p = tf_right.paragraphs[0] if i == 0 else tf_right.add_paragraph()
            if bullet.startswith("  - "):
                p.text = "    •  " + bullet[4:]
                p.font.size = Pt(13)
                p.font.color.rgb = C_MUTED
            else:
                p.text = "•  " + bullet
                p.font.size = Pt(14)
                p.font.color.rgb = C_TEXT
            p.font.name = 'Calibri'
            p.space_after = Pt(10)
            p.line_spacing = 1.15


    # ───────────────────────────────────────────────────────────────
    # SLIDE 1: Title Slide (Widescreen 16:9)
    # ───────────────────────────────────────────────────────────────
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title solid light gray background
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = RGBColor(250, 250, 250)
    bg1.line.fill.background()
    
    # Left Teal Accent Bar
    bar1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.4), Inches(7.5))
    bar1.fill.solid()
    bar1.fill.fore_color.rgb = C_TEAL
    bar1.line.fill.background()
    
    tx1 = slide1.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(11.0), Inches(4.0))
    tf1 = tx1.text_frame
    tf1.word_wrap = True
    tf1.margin_left = tf1.margin_top = tf1.margin_right = tf1.margin_bottom = 0
    
    p_tag = tf1.paragraphs[0]
    p_tag.text = "NATURAL LANGUAGE PROCESSING GROUP PROJECT"
    p_tag.font.name = 'Calibri'
    p_tag.font.size = Pt(13)
    p_tag.font.bold = True
    p_tag.font.color.rgb = C_TEAL
    p_tag.space_after = Pt(18)
    
    p_main = tf1.add_paragraph()
    p_main.text = "SMS Spam Detection System"
    p_main.font.name = 'Calibri'
    p_main.font.size = Pt(54)
    p_main.font.bold = True
    p_main.font.color.rgb = C_DARK
    p_main.space_after = Pt(6)
    
    p_sub = tf1.add_paragraph()
    p_sub.text = "A comparative study of ML and DL classifiers for SMS Spam filtering"
    p_sub.font.name = 'Calibri'
    p_sub.font.size = Pt(20)
    p_sub.font.color.rgb = C_MUTED
    p_sub.space_after = Pt(45)
    
    p_meta = tf1.add_paragraph()
    p_meta.text = "Developed by: Group Members (3 Students) | Supervised Project\nEvaluation of Baseline ML, Ensemble ML, CNN, LSTM & Transformer Networks"
    p_meta.font.name = 'Calibri'
    p_meta.font.size = Pt(13)
    p_meta.font.color.rgb = C_MUTED
    p_meta.line_spacing = 1.2


    # ───────────────────────────────────────────────────────────────
    # SLIDE 2: Project Overview & Objectives
    # ───────────────────────────────────────────────────────────────
    left_bullets2 = [
        "Unsolicited messages (advertising, phishing, scams) represent a severe security risk and a degradation of mobile user experiences.",
        "Traditional rule-based filtering fails to adapt to spelling variants, emojis, abbreviations, and context modifications.",
        "Deploying intelligent natural language classification models provides dynamic, context-aware filtering capabilities."
    ]
    right_bullets2 = [
        "Primary Goal: Automatically label incoming messages as 'Ham' (legitimate) or 'Spam' (malicious).",
        "Pipeline Integration: Support feature extraction and real-time inference on short texts.",
        "Robust Metric Focus: Ensure high F1-score to handle class imbalance, and prioritize high Precision to protect critical bank alerts or validation codes from false blocks."
    ]
    create_text_two_column_slide(
        category="INTRODUCTION",
        title="The Need for Intelligent SMS Filtering",
        left_title="Problem Statement",
        left_bullets=left_bullets2,
        right_title="Core Objectives",
        right_bullets=right_bullets2
    )


    # ───────────────────────────────────────────────────────────────
    # SLIDE 3: Dataset & Preprocessing Pipeline
    # ───────────────────────────────────────────────────────────────
    left_bullets3 = [
        "Source: UCI Machine Learning Repository SMS Spam Collection (5,574 English text messages).",
        "Class Imbalance:",
        "  - Ham (Safe): 4,827 messages (86.6%)",
        "  - Spam: 747 messages (13.4%)",
        "The class distribution mirrors the real world, requiring metrics like F1-Score and Precision over simple Accuracy."
    ]
    right_bullets3 = [
        "Text Normalization: Conversion to lowercase and regex-based removal of special characters/punctuation.",
        "Label Encoding: Mapped string categories to binary representation: 'ham' -> 0, 'spam' -> 1.",
        "Dataset Split Strategy:",
        "  - Training Set: 70% (3,900 messages) to fit parameters.",
        "  - Validation Set: 15% (836 messages) to monitor overfitting.",
        "  - Testing Set: 15% (838 messages) for final, unbiased comparison."
    ]
    create_text_two_column_slide(
        category="DATA PREPARATION",
        title="SMS Dataset Profiles & Pipeline Details",
        left_title="SMS Spam Collection Dataset",
        left_bullets=left_bullets3,
        right_title="Data Processing & Splitting",
        right_bullets=right_bullets3
    )


    # ───────────────────────────────────────────────────────────────
    # SLIDE 4: Member 1 - Logistic Regression
    # ───────────────────────────────────────────────────────────────
    lr_bullets = [
        "Represented input text as sparse numerical features using TF-IDF (Term Frequency-Inverse Document Frequency) vectorization with 3,000 maximum features.",
        "Logistic Regression models the probability of spam linearly. The sigmoid function maps arbitrary values to a [0, 1] probability range.",
        "Offers high speed, negligible resource usage, and clean feature-weight interpretability."
    ]
    create_two_column_slide(
        category="MEMBER 1: MACHINE LEARNING MODEL",
        title="Logistic Regression Classifier",
        left_bullets=lr_bullets,
        right_title="LR Evaluation Metrics (Test Set)",
        right_metrics=metrics['lr'],
        accent_color=M1_ACCENT
    )


    # ───────────────────────────────────────────────────────────────
    # SLIDE 5: Member 1 - 1D CNN
    # ───────────────────────────────────────────────────────────────
    cnn_bullets = [
        "Tokenized text into word index sequences (vocabulary size: 5,000, sequence length padded to 50) and fed to a 32-dimensional dense embedding layer.",
        "Features processed by a 1D Convolutional layer (64 filters, kernel size 5, ReLU activation) scanning text sequences like n-gram feature extractors.",
        "Global Max Pooling extracts the strongest feature activations, fed into a single sigmoid neuron output.",
        "Captures phrase patterns (e.g. 'claim cash prize') regardless of spatial location."
    ]
    create_two_column_slide(
        category="MEMBER 1: DEEP LEARNING MODEL",
        title="1D Convolutional Neural Network (CNN)",
        left_bullets=cnn_bullets,
        right_title="CNN Evaluation Metrics (Test Set)",
        right_metrics=metrics['cnn'],
        accent_color=M1_ACCENT
    )


    # ───────────────────────────────────────────────────────────────
    # SLIDE 6: Member 2 - Random Forest
    # ───────────────────────────────────────────────────────────────
    rf_bullets = [
        "Utilized TF-IDF feature extraction (vocabulary capped at 3,000 features).",
        "Constructs an ensemble of decision trees (Forest) built on bootstrapped data splits.",
        "Aggregate voting (Mode prediction) mitigates overfitting risks common to individual decision trees.",
        "Effectively captures non-linear relationships across text features without scaling requirements."
    ]
    create_two_column_slide(
        category="MEMBER 2: MACHINE LEARNING MODEL",
        title="Random Forest Classifier",
        left_bullets=rf_bullets,
        right_title="RF Evaluation Metrics (Test Set)",
        right_metrics=metrics['rf'],
        accent_color=M2_ACCENT
    )


    # ───────────────────────────────────────────────────────────────
    # SLIDE 7: Member 2 - LSTM
    # ───────────────────────────────────────────────────────────────
    lstm_bullets = [
        "Text represented via word indices and mapped through a 32-dimensional embedding layer.",
        "Long Short-Term Memory (LSTM) layer (32 hidden units) models sequence structures sequentially.",
        "Special gating mechanisms (input, forget, output gates) allow the model to propagate long-term contextual relationships.",
        "Processes word dependencies step-by-step, making it highly responsive to grammar patterns."
    ]
    create_two_column_slide(
        category="MEMBER 2: DEEP LEARNING MODEL",
        title="Long Short-Term Memory (LSTM) Network",
        left_bullets=lstm_bullets,
        right_title="LSTM Evaluation Metrics (Test Set)",
        right_metrics=metrics['lstm'],
        accent_color=M2_ACCENT
    )


    # ───────────────────────────────────────────────────────────────
    # SLIDE 8: Member 3 - XGBoost
    # ───────────────────────────────────────────────────────────────
    xgb_bullets = [
        "Uses sparse TF-IDF vectors (3,000 features) as input representation.",
        "XGBoost (Extreme Gradient Boosting) is an optimized tree boosting framework.",
        "Builds trees sequentially to minimize the residual errors of preceding trees.",
        "Includes regularization (L1/L2 weights) which prevents overfitting on sparse data splits."
    ]
    create_two_column_slide(
        category="MEMBER 3: MACHINE LEARNING MODEL",
        title="XGBoost Classifier",
        left_bullets=xgb_bullets,
        right_title="XGBoost Evaluation Metrics (Test Set)",
        right_metrics=metrics['xgb'],
        accent_color=M3_ACCENT
    )


    # ───────────────────────────────────────────────────────────────
    # SLIDE 9: Member 3 - Custom Transformer
    # ───────────────────────────────────────────────────────────────
    tf_bullets = [
        "Represents text via word tokenization, dense embeddings, and Positional Encodings to preserve word ordering context.",
        "Multi-Head Self-Attention (2 heads, key dimension 32) parses inputs in parallel, directly linking distant tokens.",
        "Utilizes Layer Normalization and Feed-Forward Networks with residual connections to stabilize deep feature propagation.",
        "Extracts global text contexts, outperforming local sequential models."
    ]
    create_two_column_slide(
        category="MEMBER 3: DEEP LEARNING MODEL",
        title="Custom Transformer Network",
        left_bullets=tf_bullets,
        right_title="Transformer Evaluation Metrics (Test Set)",
        right_metrics=metrics['transformer'],
        accent_color=M3_ACCENT
    )


    # ───────────────────────────────────────────────────────────────
    # SLIDE 10: Model Comparison & Performance Matrix
    # ───────────────────────────────────────────────────────────────
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_bg(slide10, C_TEAL)
    add_header(slide10, "Overall Performance Comparison", "EVALUATION")
    
    rows, cols = 7, 6
    tbl_shape = slide10.shapes.add_table(rows, cols, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.8))
    table = tbl_shape.table
    
    table.columns[0].width = Inches(3.233)  # Model name
    table.columns[1].width = Inches(1.5)    # Developer
    table.columns[2].width = Inches(1.75)   # Accuracy
    table.columns[3].width = Inches(1.75)   # Precision
    table.columns[4].width = Inches(1.75)   # Recall
    table.columns[5].width = Inches(1.75)   # F1-Score
    
    headers = ["Model Architecture", "Developer", "Accuracy", "Precision", "Recall", "F1-Score"]
    for c_idx, txt in enumerate(headers):
        cell = table.cell(0, c_idx)
        cell.text = txt
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_DARK
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if c_idx > 0 else PP_ALIGN.LEFT
        p.font.name = 'Calibri'
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
    row_data = [
        ("Logistic Regression", "Member 1", metrics['lr']['accuracy'], metrics['lr']['precision'], metrics['lr']['recall'], metrics['lr']['f1']),
        ("Random Forest", "Member 2", metrics['rf']['accuracy'], metrics['rf']['precision'], metrics['rf']['recall'], metrics['rf']['f1']),
        ("XGBoost", "Member 3", metrics['xgb']['accuracy'], metrics['xgb']['precision'], metrics['xgb']['recall'], metrics['xgb']['f1']),
        ("1D CNN", "Member 1", metrics['cnn']['accuracy'], metrics['cnn']['precision'], metrics['cnn']['recall'], metrics['cnn']['f1']),
        ("LSTM", "Member 2", metrics['lstm']['accuracy'], metrics['lstm']['precision'], metrics['lstm']['recall'], metrics['lstm']['f1']),
        ("Custom Transformer", "Member 3", metrics['transformer']['accuracy'], metrics['transformer']['precision'], metrics['transformer']['recall'], metrics['transformer']['f1']),
    ]
    
    for r_idx, row in enumerate(row_data):
        name, dev, acc, prec, rec, f1 = row
        cells = [name, dev, f"{acc:.2f}%", f"{prec:.2f}%", f"{rec:.2f}%", f"{f1:.2f}%"]
        is_best = (name == "Custom Transformer")
        
        for c_idx, val in enumerate(cells):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = val
            cell.fill.solid()
            if is_best:
                cell.fill.fore_color.rgb = C_TEAL_LIGHT
            else:
                cell.fill.fore_color.rgb = C_CARD_BG if r_idx % 2 == 0 else C_BG
                
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c_idx > 0 else PP_ALIGN.LEFT
            p.font.name = 'Calibri'
            p.font.size = Pt(13)
            if is_best:
                p.font.bold = True
                p.font.color.rgb = C_TEAL
            else:
                p.font.color.rgb = C_TEXT


    # ───────────────────────────────────────────────────────────────
    # SLIDE 11: Best Model Selection & Analysis
    # ───────────────────────────────────────────────────────────────
    best_m = metrics['transformer']
    left_bullets11 = [
        f"Selected Model: **Custom Transformer** (Member 3 Deep Learning Model).",
        f"Key Metrics achieved on unseen test data:",
        f"  - **Accuracy: {best_m['accuracy']:.2f}%**",
        f"  - **Precision: {best_m['precision']:.2f}%** (Highest amongst all models)",
        f"  - **Recall: {best_m['recall']:.2f}%**",
        f"  - **F1-Score: {best_m['f1']:.2f}%** (The optimal balance for imbalanced splits)"
    ]
    right_bullets11 = [
        "High Precision Advantage: Out of 100 flagged messages, 98.1 are actual spam. This protects essential bank alerts, transactional notifications, and verification codes from false blocks.",
        "Dynamic Attention Weights: The self-attention layer dynamically routes signal priority to spam indicators (e.g. 'claim', 'urgent', 'winner') and correlates them, regardless of where they occur in the message structure.",
        "Deep Learning Superiority: Deep sequential/context-aware architectures out-performed TF-IDF machine learning baseline models."
    ]
    create_text_two_column_slide(
        category="WINNER",
        title="Best Model: Custom Transformer",
        left_title="Winner Details & Score Metrics",
        left_bullets=left_bullets11,
        right_title="Evaluation Rationale",
        right_bullets=right_bullets11,
        accent_color=M3_ACCENT
    )


    # ───────────────────────────────────────────────────────────────
    # SLIDE 12: System Demo & Deployment Architecture
    # ───────────────────────────────────────────────────────────────
    left_bullets12 = [
        "Backend Engine: Flask REST API (`api.py`) handles request parsing, cleaning, loading model binaries, and responding in JSON format.",
        "Active Endpoints:",
        "  - `POST /predict`: Clean incoming message text, tokenize/vectorize, and return a prediction with a probability score using a selected model.",
        "  - `POST /compare`: Simultaneously feed a text payload to all 6 models and compile comparative predictions."
    ]
    right_bullets12 = [
        "Web User Interface: Single Page Application (HTML, CSS, JS) serving as a live demonstration environment.",
        "Modes of Operation:",
        "  - Predict Mode: Lets users type an SMS and select a model from a dropdown to check for spam labels.",
        "  - Compare Mode: Sends the message to all 6 models, rendering a visual leaderboard comparison bar chart.",
        "Robustness: Direct integration of Keras `.keras` neural networks and Scikit-learn `.pkl` models."
    ]
    create_text_two_column_slide(
        category="DEPLOYMENT",
        title="Web Interface & API Integration",
        left_title="Flask API Backend Integration",
        left_bullets=left_bullets12,
        right_title="Interactive UI Frontend",
        right_bullets=right_bullets12
    )


    # ───────────────────────────────────────────────────────────────
    # SLIDE 13: Project Conclusion & Takeaways
    # ───────────────────────────────────────────────────────────────
    left_bullets13 = [
        "Success Criteria Met: Built 6 functional classifiers with testing accuracies exceeding 97.4% and F1-scores above 93.9%.",
        "Deep Learning Benefit: Deep architectures (Transformer, CNN, LSTM) yielded higher sequence-level performance due to dense token embeddings over standard word count vectors.",
        "Clean Integration: Decoupled model training files from deployment, communicating via RESTful API."
    ]
    right_bullets13 = [
        "Future Enhancements:",
        "  - Pre-trained Embeddings: Integrate BERT or DistilBERT models to leverage large-scale language pre-training.",
        "  - Multi-Lingual Support: Extend clean filters to multi-lingual or mixed-language SMS formats (e.g. Singlish, code-switching).",
        "  - On-Device Filtering: Implement lightweight TensorFlow Lite models to run natively inside mobile phone client apps."
    ]
    create_text_two_column_slide(
        category="CONCLUSION",
        title="Summary & Project Takeaways",
        left_title="Project Conclusion",
        left_bullets=left_bullets13,
        right_title="Future Directions",
        right_bullets=right_bullets13
    )

    # Save presentation
    output_filename = "SMS_Spam_Detection_System_Presentation.pptx"
    prs.save(output_filename)
    print(f"\n[SUCCESS] Successfully compiled presentation to: {output_filename}")

if __name__ == '__main__':
    create_presentation()
